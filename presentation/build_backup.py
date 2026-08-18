"""Backup / appendix deck for Q&A defense (separate file)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
INK = RGBColor(0x1A, 0x18, 0x30); MUT = RGBColor(0x55, 0x53, 0x6B)
ACC = RGBColor(0x6A, 0x5A, 0xF9); RED = RGBColor(0xC0, 0x39, 0x2B)
LGT = RGBColor(0xF3, 0xF2, 0xFB); WHT = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Segoe UI"; W, H = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide(bg=WHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, W, H); r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor; return tf


def para(tf, text, size, color=INK, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, space=8, lh=1.1):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space); p.line_spacing = lh
    for seg in (text if isinstance(text, list) else [(text, {})]):
        t, o = seg if isinstance(seg, tuple) else (seg, {})
        r = p.add_run(); r.text = t; f = r.font
        f.size = Pt(o.get("size", size)); f.name = FONT
        f.bold = o.get("bold", bold); f.italic = o.get("italic", italic)
        f.color.rgb = o.get("color", color)
    return p


def head(s, eyebrow, title, n):
    tf = box(s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.5))
    para(tf, eyebrow.upper(), 14, ACC, bold=True, first=True)
    tf = box(s, Inches(0.7), Inches(0.9), Inches(11.9), Inches(0.9))
    para(tf, title, 30, INK, bold=True, first=True, lh=1.0)
    tf = box(s, Inches(0.7), Inches(7.02), Inches(12), Inches(0.4))
    para(tf, [("CALM-Net  ·  Appendix (backup)", {"size": 11, "color": MUT}),
              (f"          B{n}", {"size": 11, "color": MUT})], 11, first=True)


def bullets(s, x, y, w, items, size=18, gap=10):
    tf = box(s, x, y, w, Inches(4.5))
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            para(tf, [("• ", {"size": size, "color": ACC, "bold": True}),
                      (it[0], {"size": size, "bold": True, "color": INK}),
                      (it[1], {"size": size, "color": MUT})], size, first=(i == 0), space=gap, lh=1.15)
        else:
            para(tf, [("• ", {"size": size, "color": ACC, "bold": True}),
                      (it, {"size": size, "color": INK})], size, first=(i == 0), space=gap, lh=1.15)
    return tf


# ---- B0 divider ----
s = slide(INK)
bar = s.shapes.add_shape(1, 0, Inches(3.55), W, Inches(0.05)); bar.fill.solid()
bar.fill.fore_color.rgb = ACC; bar.line.fill.background(); bar.shadow.inherit = False
tf = box(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.2))
para(tf, "Appendix", 48, WHT, bold=True, first=True)
tf = box(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.8))
para(tf, "Backup slides for Q&A — evidence, method detail, evaluation plan", 20,
     RGBColor(0xCE, 0xC9, 0xF7), first=True)

# ---- B1 confound across 7 ----
s = slide(); head(s, "Evidence · the confound is not an average", "A motion sensor beats EEG on every subject", 1)
s.shapes.add_picture(str(HERE / "bk_confound7.png"), Inches(0.7), Inches(2.1), width=Inches(8.1))
tf = box(s, Inches(9.1), Inches(2.3), Inches(3.9), Inches(4.2), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("IMU-only mean ", {"size": 18, "color": INK}), ("0.92", {"size": 18, "bold": True, "color": RED}),
          (" vs EEG ", {"size": 18, "color": INK}), ("0.84", {"size": 18, "bold": True, "color": ACC})],
     18, first=True, space=14, lh=1.2)
para(tf, [("Perfect (", {"size": 17, "color": MUT}), ("1.00", {"size": 17, "bold": True, "color": RED}),
          (") for S4 and S7 — you cannot get 1.00 on a biological signal; that is the label read twice.",
           {"size": 17, "color": MUT})], 17, space=14, lh=1.2)
para(tf, "The walk/stop label is a function of the exoskeleton's motion.", 17, INK, bold=True, lh=1.2)

# ---- B2 apples-to-apples ----
s = slide(); head(s, "Fair comparison · are we behind EEGNet?", "No — the gap was the confound, not the model", 2)
s.shapes.add_picture(str(HERE / "bk_apples.png"), Inches(0.7), Inches(2.1), width=Inches(7.3))
tf = box(s, Inches(8.3), Inches(2.3), Inches(4.6), Inches(4.2), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "When EEGNet is put through the SAME movement-invariant test, it also falls to 0.65.",
     18, INK, bold=True, first=True, space=14, lh=1.2)
para(tf, "The famous 13-point “gap” to EEGNet is entirely the confound: inflated-EEGNet vs "
     "invariant-us. On equal footing, honest decoding is tied.", 17, MUT, space=14, lh=1.2)
para(tf, "Our edge is elsewhere: calibration, coverage guarantee, and safety.", 17, ACC, bold=True, lh=1.2)

# ---- B3 architecture detail ----
s = slide(); head(s, "Proposed method · detail", "CALM-Net, four stages", 3)
pic = s.shapes.add_picture(str(HERE.parent / "paper" / "fig_architecture.png"),
                           Inches(2), Inches(1.9), height=Inches(2.85))
pic.left = int((W - pic.width) / 2)
bullets(s, Inches(0.7), Inches(5.0), Inches(12),
        [("Encoder — ", "multi-band spatial features from EEG + head/exo IMU + EOG."),
         ("MID — ", "adversarial + HSIC penalty force the intent code to be independent of the IMU."),
         ("LSC — ", "temperature scaling + adaptive conformal → distribution-free coverage under drift."),
         ("SAS — ", "commit only if confident, unambiguous, and uncontaminated; else safe STOP.")],
        size=16, gap=6)

# ---- B4 method: MID ----
s = slide(); head(s, "Method · motion-invariant disentanglement", "How we remove movement from the neural code", 4)
bullets(s, Inches(0.7), Inches(2.15), Inches(12), [
    ("Target φ(M): ", "a 12-D descriptor of head + exo IMU (accel/gyro magnitude, mean/std/max)."),
    ("Split: ", "the representation is divided into an intent subspace and an artefact subspace."),
    ("Adversary (gradient reversal): ", "the encoder is trained so the intent code CANNOT predict φ(M)."),
    ("HSIC penalty: ", "a nonlinear independence term drives intent ⟂ movement (beyond linear decorrelation)."),
    ("Artefact head: ", "the other subspace is trained TO predict φ(M), absorbing the movement."),
    ("Result: ", "the classifier reads only the intent subspace → provably movement-invariant decisions."),
], size=18, gap=11)

# ---- B5 method: LSC + evaluation ----
s = slide(); head(s, "Method · guarantee + evaluation plan", "How we stay calibrated, and how we test it", 5)
tf = box(s, Inches(0.7), Inches(2.1), Inches(6.0), Inches(4.5))
para(tf, "Longitudinal self-calibration", 18, ACC, bold=True, first=True, space=8)
for it in [("Temperature scaling", " for honest confidence on a held-out split."),
           ("Split + adaptive conformal", " updates the threshold online: q ← q + η(α − err)."),
           ("Guarantee", " prediction-set coverage tracks its target as the signal drifts.")]:
    para(tf, [("• ", {"size": 16, "color": ACC, "bold": True}), (it[0], {"size": 16, "bold": True}),
              (it[1], {"size": 16, "color": MUT})], 16, space=8, lh=1.15)
tf = box(s, Inches(7.0), Inches(2.1), Inches(5.9), Inches(4.5))
para(tf, "Evaluation protocol", 18, ACC, bold=True, first=True, space=8)
for it in [("Longitudinal split", " — train early sessions, test later ones (no leakage)."),
           ("Balanced accuracy + ECE + Brier", " — accuracy and calibration."),
           ("Coverage @ target + wrong-walk rate", " — the safety metrics."),
           ("Movement-invariance probe", " — intent→IMU R² confirms honesty."),
           ("Baselines", " — EEGNet, EEG-Conformer, and the IMU-only control.")]:
    para(tf, [("• ", {"size": 16, "color": ACC, "bold": True}), (it[0], {"size": 16, "bold": True}),
              (it[1], {"size": 16, "color": MUT})], 16, space=8, lh=1.15)

# ---- B6 dataset + references ----
s = slide(); head(s, "Dataset & key references", "NeuroRex (OpenNeuro ds007788)", 6)
tf = box(s, Inches(0.7), Inches(2.1), Inches(12), Inches(1.5))
para(tf, [("7 healthy adults × 9 sessions", {"size": 18, "bold": True, "color": INK}),
          (" over 14–80 days · 60-ch EEG + 4 EOG @ 100 Hz · head & exoskeleton IMUs · "
           "exoskeleton feedback logs · Walk vs Stop motor imagery · BIDS, CC0.",
           {"size": 18, "color": MUT})], 18, first=True, lh=1.2)
tf = box(s, Inches(0.7), Inches(3.5), Inches(12), Inches(3.2))
para(tf, "Selected references", 15, ACC, bold=True, first=True, space=8)
for r in ["Sarkar et al. (2026) — NeuroRex dataset. Scientific Data.",
          "Lawhern et al. (2018) — EEGNet. J. Neural Eng.",
          "Song et al. (2023) — EEG Conformer. IEEE TNSRE.",
          "Ganin et al. (2016) — Domain-adversarial training (gradient reversal). JMLR.",
          "Gibbs & Candès (2021) — Adaptive conformal inference. NeurIPS.",
          "Geifman & El-Yaniv (2019) — SelectiveNet. ICML."]:
    para(tf, r, 15, MUT, space=6, lh=1.15)

# ---- B7 architecture study: one ceiling ----
s = slide(); head(s, "Result · architecture study", "We tried 5 architectures — none beats the ceiling", 7)
s.shapes.add_picture(str(HERE / "bk_ceiling.png"), Inches(0.7), Inches(2.1), width=Inches(7.7))
tf = box(s, Inches(8.7), Inches(2.3), Inches(4.3), Inches(4.2), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "EEGNet, band-power, Riemannian + attention, and a graph net over the scalp geometry — "
     "all land at ~0.65 when made movement-invariant.", 17, INK, bold=True, first=True, space=14, lh=1.2)
para(tf, "The bottleneck is not the model. The movement-invariant neural signal is simply weak — "
     "and equal for every architecture.", 17, MUT, lh=1.2)

# ---- B8 exhaustive investigation ----
s = slide(); head(s, "Result · we tried to raise the ceiling", "Five independent attempts — all confirm it", 8)
bullets(s, Inches(0.7), Inches(2.15), Inches(12), [
    ("Stronger disentanglement (HSIC): ", "richer model’s accuracy collapses to 0.65 — its lead was leakage."),
    ("Pre-movement window: ", "the IMU still scores 0.94 at the transition — no movement-free window exists."),
    ("EEG + IMU fusion: ", "equals IMU-only (0.92); EEG adds nothing over the motion sensor."),
    ("5× more training data: ", "helps weak subjects reach the ceiling, but never exceeds it."),
    ("Graph over electrode geometry: ", "a new prior, still ~0.65 invariant."),
], size=18, gap=13)
bar = s.shapes.add_shape(1, Inches(0.7), Inches(5.9), Inches(11.93), Inches(0.95))
bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background(); bar.shadow.inherit = False
bt = bar.text_frame; bt.word_wrap = True; bt.vertical_anchor = MSO_ANCHOR.MIDDLE
bt.margin_left = Pt(18); bt.margin_right = Pt(18)
para(bt, "The ceiling is real: this is a movement benchmark, not a brain-decoding benchmark. "
     "That is the finding.", 18, WHT, bold=True, first=True, lh=1.15)

# ---- B9 framework wins ----
s = slide(); head(s, "Result · what CALM-Net actually delivers", "The win is trustworthiness, not accuracy", 9)
s.shapes.add_picture(str(HERE.parent / "paper" / "fig_full.png"), Inches(0.7), Inches(2.15), width=Inches(7.7))
tf = box(s, Inches(8.7), Inches(2.2), Inches(4.3), Inches(4.4), anchor=MSO_ANCHOR.MIDDLE)
for t, d in [("Coverage guarantee", " held on all 7 subjects (0.90) as the signal drifts — static conformal fails."),
             ("Calibration", " improved across the session gap (ECE 0.15 → 0.11)."),
             ("Wrong-walk error", " cut 3× (0.14 → 0.05) — the dangerous error, controlled.")]:
    para(tf, [("✓  ", {"size": 18, "bold": True, "color": ACC}), (t, {"size": 18, "bold": True}),
              (d, {"size": 18, "color": MUT})], 18, first=(t == "Coverage guarantee"), space=14, lh=1.18)

out = HERE / "CALM-Net_backup_appendix.pptx"
prs.save(str(out))
print("saved", out.name, "|", len(prs.slides._sldIdLst), "slides")
