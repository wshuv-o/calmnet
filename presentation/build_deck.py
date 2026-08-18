"""Midterm proposal deck (11 slides), old-school academic style.
Serif type, black on white, no accent colour, plain ruled tables."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
BLK = RGBColor(0x11, 0x11, 0x11)      # text
GRY = RGBColor(0x55, 0x55, 0x55)      # secondary text
LINE = RGBColor(0x00, 0x00, 0x00)     # rules
BORD = RGBColor(0x44, 0x44, 0x44)     # box outlines
WHT = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Times New Roman"
W, H = Inches(13.333), Inches(7.5)
NSL = 11
prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide(bg=WHT):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, W, H); r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor; return tf


def para(tf, text, size, color=BLK, bold=False, italic=False, first=False,
         align=PP_ALIGN.LEFT, space=8, lh=1.1):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space); p.line_spacing = lh
    for seg in (text if isinstance(text, list) else [(text, {})]):
        t, o = seg if isinstance(seg, tuple) else (seg, {})
        r = p.add_run(); r.text = t; f = r.font
        f.size = Pt(o.get("size", size)); f.name = FONT
        f.bold = o.get("bold", bold); f.italic = o.get("italic", italic)
        f.color.rgb = o.get("color", color)
    return p


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(1, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def obox(s, x, y, w, h, lw=1.0):
    """White box with a thin outline (no fill colour)."""
    return rect(s, x, y, w, h, fill=WHT, line=BORD, lw=lw)


def hrule(s, x, y, w, pt=1.0):
    rect(s, x, y, w, Pt(pt), fill=LINE)


def header(s, text):
    tf = box(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.85))
    para(tf, text, 30, BLK, bold=True, first=True, lh=1.0)
    hrule(s, Inches(0.7), Inches(1.32), Inches(11.93), pt=1.4)


def footer(s, n):
    return  # no running footer (old-school, kept as no-op)


def no_style(gt):
    """Strip the theme table style so cells render plain white."""
    tblPr = gt._tbl.find(qn('a:tblPr'))
    if tblPr is None:
        return
    for a in ('firstRow', 'firstCol', 'lastRow', 'lastCol', 'bandRow', 'bandCol'):
        if a in tblPr.attrib:
            tblPr.set(a, '0')
    sid = tblPr.find(qn('a:tableStyleId'))
    if sid is None:
        sid = tblPr.makeelement(qn('a:tableStyleId'), {})
        tblPr.append(sid)
    sid.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'  # No Style, No Grid


def tcell(gt, ri, ci, text, size, bold, align):
    cell = gt.cell(ri, ci); cell.fill.background()
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(10); cell.margin_right = Pt(6); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; f = r.font; f.name = FONT
    f.size = Pt(size); f.bold = bold; f.color.rgb = BLK


# ═══ 1 · TITLE ═══
s = slide()
tf = box(s, Inches(0.5), Inches(2.15), Inches(12.33), Inches(1.2))
para(tf, "CALM-Net", 54, BLK, bold=True, first=True, align=PP_ALIGN.CENTER)
hrule(s, Inches(4.42), Inches(3.32), Inches(4.5), pt=1.5)
tf = box(s, Inches(1.4), Inches(3.55), Inches(10.53), Inches(1.4))
para(tf, "A Calibrated, Abstaining, Longitudinal EEG Decoder for Safe Lower-Limb Exoskeleton Control",
     24, BLK, first=True, align=PP_ALIGN.CENTER, lh=1.15)
tf = box(s, Inches(1.4), Inches(5.95), Inches(10.53), Inches(1.0))
para(tf, "Md Wahiduzzaman Suva  [26-94088-2]", 15, GRY, first=True, align=PP_ALIGN.CENTER, space=5)
para(tf, "Esme Moula Chowdhury Abha  [26-94089-2]", 15, GRY, align=PP_ALIGN.CENTER)

# ═══ 2 · PROBLEM STATEMENT ═══
s = slide(); header(s, "Problem Statement")
tf = box(s, Inches(0.7), Inches(1.75), Inches(12), Inches(1.6))
para(tf, "For someone who cannot walk after a spinal cord injury or stroke, an exoskeleton driven by "
     "their own brain activity restores mobility. Only non-invasive EEG requires no surgery, so it is "
     "the only version most patients could ever use.", 20, BLK, first=True, lh=1.25)
q = obox(s, Inches(0.7), Inches(3.75), Inches(11.93), Inches(1.35), lw=1.3)
qt = q.text_frame; qt.word_wrap = True; qt.vertical_anchor = MSO_ANCHOR.MIDDLE; qt.margin_left = Pt(18)
para(qt, "A confident wrong “walk” command at the top of a staircase is a fall, not a rounding error.",
     22, BLK, bold=True, first=True, lh=1.15)
tf = box(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.1))
para(tf, [("Accuracy is the wrong target. ", {"size": 19, "bold": True, "color": BLK}),
          ("A machine that moves your legs needs honest confidence and the ability to stop.",
           {"size": 19, "color": GRY})], 19, first=True, lh=1.2)

# ═══ 3 · KEY CHALLENGES ═══
s = slide(); header(s, "Key Challenges")
cards = [("1", "What is actually decoded?", "The exoskeleton moves the wearer during walking, so accuracy can reflect movement rather than neural intent."),
         ("2", "The signal drifts", "EEG is non-stationary; a decoder calibrated one day is miscalibrated the next, over weeks of use."),
         ("3", "No reject option", "Deployed decoders must guess every input; they cannot abstain to a safe stop when unsure.")]
x = 0.7; cw = 3.94; gap = 0.155
for tag, t, d in cards:
    c = obox(s, Inches(x), Inches(1.85), Inches(cw), Inches(3.7))
    ct = c.text_frame; ct.word_wrap = True; ct.vertical_anchor = MSO_ANCHOR.TOP
    ct.margin_left = Pt(16); ct.margin_right = Pt(16); ct.margin_top = Pt(18)
    para(ct, tag, 30, BLK, bold=True, first=True, space=6)
    para(ct, t, 20, BLK, bold=True, space=10, lh=1.05)
    para(ct, d, 16, GRY, lh=1.22)
    x += cw + gap

# ═══ 4 · OBJECTIVES ═══
s = slide(); header(s, "Objectives")
objs = [("O1", "Movement-invariant intent", "Decode neural intent provably invariant to the exoskeleton’s own motion."),
        ("O2", "Calibrated across time", "Stay honestly calibrated across weeks of drift, without full retraining."),
        ("O3", "Abstain with a guarantee", "Reject low-confidence inputs to a safe stop, with a distribution-free guarantee.")]
y = 1.85
for tag, t, d in objs:
    c = obox(s, Inches(0.7), Inches(y), Inches(1.0), Inches(1.0))
    ctf = c.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(ctf, tag, 24, BLK, bold=True, first=True, align=PP_ALIGN.CENTER)
    tf = box(s, Inches(1.95), Inches(y - 0.03), Inches(10.6), Inches(1.1))
    para(tf, t, 21, BLK, bold=True, first=True, space=2)
    para(tf, d, 16, GRY, lh=1.12)
    y += 1.36
tf = box(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.6))
para(tf, [("Overarching question:  ", {"size": 17, "bold": True, "color": BLK}),
          ("can we trade inflated accuracy for trustworthy, verifiable closed-loop control?",
           {"size": 17, "italic": True, "color": GRY})], 17, first=True)

# ═══ 5 · LITERATURE REVIEW ═══
s = slide(); header(s, "Literature Review")
cols = [("Lower-limb exoskeleton BCI", ["NeuroRex dataset  [Sarkar 2026]", "Deep-learning control  [Ferrero 2024]",
                                        "Exo error potentials  [Soriano 2025]", "Robot-assisted gait  [Tortora 2023]"]),
        ("Motor-imagery decoders", ["EEG-Conformer  [Song 2023]", "EEGNet  [Lawhern 2018]",
                                    "Riemannian geom.  [Barachant 2012]"]),
        ("Calibration, uncertainty", ["Adaptive conformal  [Gibbs 2021]", "Conformal pred.  [Angelopoulos 2023]",
                                      "EEG uncertainty  [Tveter 2024]"]),
        ("Drift and abstention", ["Deep adaptation  [Zhang 2023]", "Reject-option  [Hendrickx 2024]",
                                  "Selective class.  [Geifman 2017]"])]
x = 0.7; cw = 2.9; gap = 0.113
for t, items in cols:
    c = obox(s, Inches(x), Inches(1.65), Inches(cw), Inches(2.85))
    ctf = c.text_frame; ctf.word_wrap = True; ctf.vertical_anchor = MSO_ANCHOR.TOP
    ctf.margin_left = Pt(11); ctf.margin_right = Pt(9); ctf.margin_top = Pt(13)
    para(ctf, t, 15, BLK, bold=True, first=True, space=9, lh=1.05)
    for it in items:
        para(ctf, it, 12, GRY, space=7, lh=1.1)
    x += cw + gap
bar = obox(s, Inches(0.7), Inches(4.75), Inches(11.93), Inches(1.6), lw=1.6)
btf = bar.text_frame; btf.word_wrap = True; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
btf.margin_left = Pt(20); btf.margin_right = Pt(20)
para(btf, "RESEARCH GAP", 14, BLK, bold=True, first=True, space=4)
para(btf, "Rich multimodal exoskeleton datasets now exist, yet no prior work joins movement disentanglement, "
     "cross-session calibration, and safety-asymmetric abstention in one closed-loop decoder.", 19, BLK, lh=1.14)
tf = box(s, Inches(0.7), Inches(6.5), Inches(11.93), Inches(0.4))
para(tf, "Citations shown as [first author, year]; the paper cites 11 lower-limb exoskeleton BCI works and 41 references in total.",
     10.5, GRY, italic=True, first=True)

# ═══ 6 · PROPOSED METHOD: ARCHITECTURE ═══
s = slide(); header(s, "Proposed Method: CALM-Net Architecture")
pic = s.shapes.add_picture(str(HERE / "calm_arch.png"), 0, Inches(1.5), width=Inches(11.6))
pic.left = int((W - pic.width) / 2)
tf = box(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.8))
para(tf, [("A multimodal pipeline: ", {"size": 16, "bold": True, "color": BLK}),
          ("a neuro-kinematic encoder, motion-invariant disentanglement, longitudinal self-calibration, "
           "and a safety-asymmetric decision that defaults to Stop.", {"size": 16, "color": GRY})],
     16, first=True, lh=1.15)

# ═══ 7 · CORE IDEA AND HYPOTHESIS ═══
s = slide(); header(s, "Proposed Method: Core Idea and Hypothesis")
bullets = [("Record movement alongside EEG:", " the head and exoskeleton IMUs give a 12-D motion signal."),
           ("Adversarial disentanglement (MID):", " force the neural code to be statistically independent of that motion."),
           ("Decide from the invariant part only:", " the classifier never sees the movement-coupled features."),
           ("Add honesty and safety:", " temperature scaling with adaptive conformal, then abstain to Stop when unsure.")]
tf = box(s, Inches(0.7), Inches(1.75), Inches(12), Inches(3))
for i, (t, d) in enumerate(bullets):
    para(tf, [("•  ", {"size": 19, "bold": True, "color": BLK}),
              (t, {"size": 19, "bold": True, "color": BLK}), (d, {"size": 19, "color": GRY})],
         19, first=(i == 0), space=13, lh=1.18)
hb = obox(s, Inches(0.7), Inches(5.35), Inches(11.93), Inches(1.45), lw=1.3)
htf = hb.text_frame; htf.word_wrap = True; htf.vertical_anchor = MSO_ANCHOR.MIDDLE
htf.margin_left = Pt(18); htf.margin_right = Pt(18)
para(htf, [("HYPOTHESIS   ", {"size": 14, "bold": True, "color": BLK}),
           ("Disentangling movement, calibrating across sessions, and abstaining under uncertainty yields "
            "an honest, safe decoder that trades inflated accuracy for trustworthy control.",
            {"size": 18, "color": BLK, "italic": True})], 18, first=True, lh=1.15)

# ═══ 8 · PRELIMINARY ANALYSIS ═══
s = slide(); header(s, "Preliminary Analysis: Movement Confound")
s.shapes.add_picture(str(HERE / "hook_bw.png"), Inches(0.55), Inches(2.05), width=Inches(6.3))
tf = box(s, Inches(7.15), Inches(1.9), Inches(5.7), Inches(4.6), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "The exoskeleton physically moves the wearer during walk.", 18, BLK, bold=True, first=True, space=12, lh=1.15)
para(tf, [("A motion sensor with ", {"size": 17, "color": GRY}),
          ("no brain data", {"size": 17, "bold": True, "color": BLK}),
          (" scores ", {"size": 17, "color": GRY}), ("0.92", {"size": 17, "bold": True, "color": BLK}),
          (", above the best EEG decoder (0.84).", {"size": 17, "color": GRY})], 17, space=12, lh=1.18)
para(tf, "So the reported accuracy is largely the confound. This is exactly why we decode a "
     "movement-invariant signal.", 18, BLK, bold=True, lh=1.18)

# ═══ 9 · PRELIMINARY RESULTS (architecture comparison) ═══
s = slide(); header(s, "Preliminary Results: Architectures Compared")
tf = box(s, Inches(0.7), Inches(1.48), Inches(12), Inches(0.4))
para(tf, "Balanced accuracy, mean over 7 subjects, identical longitudinal protocol. "
     "For reference, an IMU sensor with no brain data scores 0.84.", 13.5, GRY, italic=True, first=True)
models = ["", "EEGNet", "Band-power", "Riemannian", "Graph net", "Conformer"]
data = [("Parameters", ["4.8 k", "4.0 k", "—", "8.4 k", "341 k"]),
        ("Standard protocol", ["0.81", "0.81", "0.81", "0.84", "0.83"]),
        ("Movement-invariant", ["0.65", "0.65", "0.65", "0.69", "—"]),
        ("Motion recoverable  R²", ["−0.38", "−0.31", "−0.47", "−0.04", "—"])]
tx, ty, tw = 0.7, 2.05, 11.93; hh, dr = 0.52, 0.52
gt = s.shapes.add_table(1 + len(data), 6, Inches(tx), Inches(ty), Inches(tw), Inches(hh + len(data) * dr)).table
no_style(gt)
gt.columns[0].width = Inches(3.03)
for ci in range(1, 6):
    gt.columns[ci].width = Inches(1.78)
gt.rows[0].height = Inches(hh)
for ci, name in enumerate(models):
    tcell(gt, 0, ci, name, 15, True, PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
for ri, (label, vals) in enumerate(data, start=1):
    gt.rows[ri].height = Inches(dr)
    tcell(gt, ri, 0, label, 14.5, True, PP_ALIGN.LEFT)
    inv = (label == "Movement-invariant")
    for ci, v in enumerate(vals, start=1):
        tcell(gt, ri, ci, v, 15, inv, PP_ALIGN.CENTER)
hrule(s, Inches(tx), Inches(ty), Inches(tw), pt=1.5)                     # top rule
hrule(s, Inches(tx), Inches(ty + hh), Inches(tw), pt=1.0)                # header rule
hrule(s, Inches(tx), Inches(ty + hh + len(data) * dr), Inches(tw), pt=1.5)  # bottom rule
tf = box(s, Inches(0.7), Inches(4.98), Inches(11.93), Inches(0.35))
para(tf, "Motion R² below zero means the exoskeleton’s motion is no longer recoverable from the neural code.",
     11.5, GRY, italic=True, first=True)
tf = box(s, Inches(0.7), Inches(5.42), Inches(11.93), Inches(1.4))
para(tf, [("Takeaway:  ", {"size": 16, "bold": True, "color": BLK}),
          ("every neural architecture collapses to about 0.65 once movement is removed, and none clearly wins.",
           {"size": 16, "color": GRY})], 16, first=True, lh=1.2)
para(tf, [("The contribution is the framework, not raw accuracy:  ", {"size": 16, "bold": True, "color": BLK}),
          ("adaptive-conformal coverage held at 0.90 (7 / 7) and wrong-walk error fell from 0.14 to 0.05.",
           {"size": 16, "color": GRY})], 16, lh=1.2)

# ═══ 10 · PER-SUBJECT RESULTS ═══
s = slide(); header(s, "Per-Subject Results")
tf = box(s, Inches(0.7), Inches(1.48), Inches(12), Inches(0.4))
para(tf, "Full CALM-Net pipeline, longitudinal (train on sessions 1 to 3, test on the later 6). "
     "Coverage target 0.90.", 13.5, GRY, italic=True, first=True)
psub = ["Subject", "Span (days)", "IMU-only", "Invariant acc.", "Coverage"]
prows = [("S01", "21", "0.86", "0.79", "0.90"),
         ("S02", "15", "0.79", "0.52", "0.90"),
         ("S03", "14", "0.98", "0.77", "0.91"),
         ("S04", "35", "0.96", "0.70", "0.90"),
         ("S05", "44", "0.89", "0.66", "0.90"),
         ("S06", "80", "0.56", "0.70", "0.91"),
         ("S07", "39", "0.83", "0.73", "0.90"),
         ("Mean", "", "0.84", "0.69", "0.90")]
tx, ty, tw = 0.7, 2.0, 11.93; hh, dr = 0.5, 0.44
gt = s.shapes.add_table(1 + len(prows), 5, Inches(tx), Inches(ty), Inches(tw), Inches(hh + len(prows) * dr)).table
no_style(gt)
gt.columns[0].width = Inches(2.0); gt.columns[1].width = Inches(2.25)
for ci in range(2, 5):
    gt.columns[ci].width = Inches(2.56)
gt.rows[0].height = Inches(hh)
for ci, name in enumerate(psub):
    tcell(gt, 0, ci, name, 14.5, True, PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
for ri, row in enumerate(prows, start=1):
    mean = row[0] == "Mean"
    gt.rows[ri].height = Inches(dr)
    for ci, v in enumerate(row):
        tcell(gt, ri, ci, v, 14, mean, PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
hrule(s, Inches(tx), Inches(ty), Inches(tw), pt=1.5)                                   # top
hrule(s, Inches(tx), Inches(ty + hh), Inches(tw), pt=1.0)                              # header
hrule(s, Inches(tx), Inches(ty + hh + (len(prows) - 1) * dr), Inches(tw), pt=0.75)     # above mean
hrule(s, Inches(tx), Inches(ty + hh + len(prows) * dr), Inches(tw), pt=1.5)            # bottom
tf = box(s, Inches(0.7), Inches(6.05), Inches(11.93), Inches(1.2))
para(tf, [("S03: ", {"size": 15, "bold": True, "color": BLK}),
          ("a motion sensor alone scores 0.98, yet the movement-invariant decode is 0.77; the gap is the "
           "confound being removed.", {"size": 15, "color": GRY})], 15, first=True, lh=1.2)
para(tf, [("S06: ", {"size": 15, "bold": True, "color": BLK}),
          ("motion is uninformative (0.56), yet EEG still decodes 0.70, genuine neural signal. Coverage holds "
           "near 0.90 for every subject, across spans up to 80 days.", {"size": 15, "color": GRY})], 15, lh=1.2)

# ═══ 11 · REFERENCES ═══
s = slide(); header(s, "References")
tf = box(s, Inches(0.7), Inches(1.42), Inches(12), Inches(0.35))
para(tf, "Works cited in the literature review, keyed to the [author, year] tags on slide 5.",
     12.5, GRY, italic=True, first=True)
refs = [
    ("[Sarkar 2026]", "S. Sarkar et al., “EEG-controlled exoskeleton for walking and standing: a longitudinal multimodal dataset,” Sci. Data, 2026. (OpenNeuro ds007788)"),
    ("[Ferrero 2024]", "L. Ferrero et al., “Brain-machine interface based on deep learning to control a lower-limb robotic exoskeleton,” J. NeuroEng. Rehabil., 21(48), 2024."),
    ("[Soriano 2025]", "P. Soriano-Segura, M. Ortiz et al., “Characterization of error-related potentials during exoskeleton command via deep learning,” J. NeuroEng. Rehabil., 2025."),
    ("[Tortora 2023]", "S. Tortora et al., “Cortical and muscular activity under robot-assisted gait modes during exoskeleton walking,” 2023."),
    ("[Song 2023]", "Y. Song et al., “EEG Conformer: convolutional transformer for EEG decoding and visualization,” IEEE TNSRE, 31:710-719, 2023."),
    ("[Lawhern 2018]", "V. J. Lawhern et al., “EEGNet: a compact CNN for EEG-based brain-computer interfaces,” J. Neural Eng., 15(5), 2018."),
    ("[Barachant 2012]", "A. Barachant et al., “Multiclass brain-computer interface classification by Riemannian geometry,” IEEE TBME, 59(4), 2012."),
    ("[Gibbs 2021]", "I. Gibbs and E. Candès, “Adaptive conformal inference under distribution shift,” NeurIPS, 2021."),
    ("[Angelopoulos 2023]", "A. N. Angelopoulos and S. Bates, “Conformal prediction: a gentle introduction,” Found. Trends Mach. Learn., 16(4), 2023."),
    ("[Tveter 2024]", "M. Tveter et al., “Advancing EEG prediction with deep learning and uncertainty estimation,” Brain Inform., 11(27), 2024."),
    ("[Zhang 2023]", "X. Zhang et al., “Priming cross-session motor-imagery classification with a universal deep domain-adaptation framework,” Neurocomputing, 556, 2023."),
    ("[Hendrickx 2024]", "K. Hendrickx et al., “Machine learning with a reject option: a survey,” Mach. Learn., 113(5), 2024."),
    ("[Geifman 2017]", "Y. Geifman and R. El-Yaniv, “Selective classification for deep neural networks,” NeurIPS, 2017."),
]
tfl = box(s, Inches(0.7), Inches(1.95), Inches(5.95), Inches(5.1))
for i, (tag, cit) in enumerate(refs[:7]):
    para(tfl, [(tag + "  ", {"size": 11, "bold": True, "color": BLK}), (cit, {"size": 10.5, "color": GRY})],
         10.5, first=(i == 0), space=9, lh=1.06)
tfr = box(s, Inches(6.9), Inches(1.95), Inches(5.75), Inches(5.1))
for i, (tag, cit) in enumerate(refs[7:]):
    para(tfr, [(tag + "  ", {"size": 11, "bold": True, "color": BLK}), (cit, {"size": 10.5, "color": GRY})],
         10.5, first=(i == 0), space=9, lh=1.06)

out = HERE / "CALM-Net_midterm_proposal.pptx"
prs.save(str(out))
print("saved", out.name, "|", len(prs.slides._sldIdLst), "slides")
